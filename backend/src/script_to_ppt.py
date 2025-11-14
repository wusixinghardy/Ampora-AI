import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs, title="Ampora Generated Presentation", subtitle="AI-Powered Educational Content"):
    """Create a title slide"""
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def add_image_slide(prs, image_path, script_text):
    """
    Add a slide with:
    - Image displayed prominently
    - Script in speaker notes
    """
    # Use blank layout to have full control
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add image (centered, taking up most of the slide)
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)  # Adjust as needed
    
    try:
        pic = slide.shapes.add_picture(image_path, left, top, width=width)
        
        # Center the image
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = int((slide_height - pic.height) / 2)
        
    except Exception as e:
        print(f"Warning: Could not add image {image_path}: {e}")
    
    # Add script to speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = script_text
    
    return slide

def create_presentation_from_scripts(scripts_json="scripts/image_scripts.json", 
                                     output_ppt="output/generated_presentation.pptx",
                                     title="Ampora Generated Presentation"):
    """
    Create PowerPoint from scripts JSON
    """
    # Load scripts
    with open(scripts_json, 'r') as f:
        scripts_data = json.load(f)
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    print("Creating title slide...")
    create_title_slide(prs, title=title)
    
    # Add image slides
    print(f"\nAdding {len(scripts_data)} image slides...")
    for idx, (image_name, data) in enumerate(scripts_data.items(), 1):
        image_path = data['path']
        script = data['script']
        
        print(f"[{idx}/{len(scripts_data)}] Adding slide for {image_name}")
        add_image_slide(prs, image_path, script)
    
    # Save presentation
    os.makedirs("output", exist_ok=True)
    prs.save(output_ppt)
    print(f"\n✓ Presentation saved to {output_ppt}")
    
    return output_ppt

if __name__ == "__main__":
    create_presentation_from_scripts()