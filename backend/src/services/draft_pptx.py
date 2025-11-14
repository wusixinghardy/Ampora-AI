from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import os

# Placeholder imports for Hardy’s and Sam’s modules
# (Replace with actual imports once their functions are implemented)
try:
    from backend.src.services.lecture import generate_lecture_bulletpoint
    from backend.src.services.visualization import generate_visualizations
except ImportError:
    # If Hardy or Sam's functions are not ready yet, use mock placeholders
    def generate_lecture_bulletpoint(concept):
        return [
            ["Intro to AI", "Definition of AI", "History and Applications"],
            ["Machine Learning Basics", "Supervised vs Unsupervised Learning"],
        ]

    def generate_visualizations(lecture_scripts):
        # Mock list of image file paths corresponding to each slide
        return [
            "backend/assets/ai_intro.png",
            "backend/assets/ml_basics.png",
        ]


def create_pptx(bullet_points, visualizations, output_file="generated_presentation.pptx"):
    """Combine bullet points and visualizations into a PowerPoint file."""
    prs = Presentation()

    for i, (bullets, vis_path) in enumerate(zip(bullet_points, visualizations)):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_tf = title_box.text_frame
        title_tf.text = f"Slide {i + 1}"

        # Add bullet points
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(4)
        height = Inches(4)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame

        for point in bullets:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)

        # Add visualization image
        if os.path.exists(vis_path):
            slide.shapes.add_picture(vis_path, Inches(5.5), Inches(1.5), width=Inches(4))
        else:
            print(f"Warning: visualization not found at {vis_path}")

    prs.save(output_file)
    print(f"✅ Presentation saved as {output_file}")
    return output_file


def main():
    """Main integration pipeline."""
    # Example concept to pass to Hardy's function
    concept = "Artificial Intelligence"

    # Call Hardy’s and Sam’s functions (or placeholders)
    bullet_points = generate_lecture_bulletpoint(concept)
    visualizations = generate_visualizations(bullet_points)

    # Combine and export PowerPoint
    output_path = create_pptx(bullet_points, visualizations)
    return output_path

# Or comment this out and do 'python -m backend.src.services.pptx' 
# And still run 'python backend/src/services/pptx.py'
if __name__ == "__main__":
    create_presentation("Linear Regression")
