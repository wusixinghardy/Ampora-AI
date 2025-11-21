"""
voice_generator.py - Text-to-Speech Voice Generation
Author: Robert Jarman

Generates voiceover audio files from Hardy's scripts.
Uses OpenAI TTS API or pyttsx3 for offline generation.
"""

import os
from typing import List, Tuple
from pathlib import Path

# Option 1: Use OpenAI TTS (recommended - higher quality)
try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

# Option 2: Fallback to local TTS
try:
    import pyttsx3
    PYTTSX3_AVAILABLE = True
except ImportError:
    PYTTSX3_AVAILABLE = False


class VoiceGenerator:
    """
    Generate voiceover audio from text scripts.
    """
    
    def __init__(self, use_openai: bool = True):
        """
        Initialize voice generator.
        
        Args:
            use_openai: If True, use OpenAI TTS API; otherwise use local TTS
        """
        self.use_openai = use_openai and OPENAI_AVAILABLE
        
        if self.use_openai:
            from src.config import OPENAI_API_KEY
            self.client = OpenAI(api_key=OPENAI_API_KEY)
            self.voice = "alloy"  # Options: alloy, echo, fable, onyx, nova, shimmer
            self.model = "tts-1"  # or "tts-1-hd" for higher quality
        elif PYTTSX3_AVAILABLE:
            self.engine = pyttsx3.init()
            # Configure voice properties
            self.engine.setProperty('rate', 150)  # Speed of speech
            self.engine.setProperty('volume', 0.9)
        else:
            raise RuntimeError("No TTS engine available. Install openai or pyttsx3")
    
    def generate_audio_openai(self, text: str, output_path: str) -> str:
        """
        Generate audio using OpenAI TTS API.
        
        Args:
            text: Script text to convert to speech
            output_path: Path to save MP3 file
            
        Returns:
            Path to generated audio file
        """
        response = self.client.audio.speech.create(
            model=self.model,
            voice=self.voice,
            input=text
        )
        
        response.stream_to_file(output_path)
        return output_path
    
    def generate_audio_local(self, text: str, output_path: str) -> str:
        """
        Generate audio using local pyttsx3.
        
        Args:
            text: Script text to convert to speech
            output_path: Path to save audio file
            
        Returns:
            Path to generated audio file
        """
        self.engine.save_to_file(text, output_path)
        self.engine.runAndWait()
        return output_path
    
    def generate_voiceover_for_scripts(
        self,
        scripts: List[str],
        output_dir: str = "backend/output/voiceovers"
    ) -> List[Tuple[str, float]]:
        """
        Generate voiceover audio files for all scripts.
        
        Args:
            scripts: List of script texts (one per slide)
            output_dir: Directory to save audio files
            
        Returns:
            List of tuples: (audio_file_path, duration_seconds)
        """
        os.makedirs(output_dir, exist_ok=True)
        
        audio_files = []
        
        print(f"\n🎙️  Generating voiceovers for {len(scripts)} slides...")
        
        for i, script in enumerate(scripts, 1):
            # Clean script text
            clean_script = script.strip()
            
            if not clean_script:
                print(f"   ⚠️  Slide {i}: Empty script, skipping")
                continue
            
            # Generate filename
            if self.use_openai:
                filename = f"slide_{i:02d}.mp3"
            else:
                filename = f"slide_{i:02d}.wav"
            
            output_path = os.path.join(output_dir, filename)
            
            try:
                print(f"   [{i}/{len(scripts)}] Generating audio...")
                
                if self.use_openai:
                    self.generate_audio_openai(clean_script, output_path)
                else:
                    self.generate_audio_local(clean_script, output_path)
                
                # Estimate duration (rough: 150 words per minute, 5 chars per word)
                word_count = len(clean_script) / 5
                duration = (word_count / 150) * 60  # seconds
                
                audio_files.append((output_path, duration))
                print(f"   ✅ Generated: {filename} (~{duration:.1f}s)")
                
            except Exception as e:
                print(f"   ❌ Error generating audio for slide {i}: {e}")
                continue
        
        print(f"\n✅ Generated {len(audio_files)} voiceover files")
        return audio_files


def generate_voiceovers_from_pptx(pptx_path: str, output_dir: str = "backend/output/voiceovers"):
    """
    Extract speaker notes from PowerPoint and generate voiceovers.
    
    Args:
        pptx_path: Path to PowerPoint file
        output_dir: Directory to save audio files
    """
    from pptx import Presentation
    
    prs = Presentation(pptx_path)
    scripts = []
    
    # Extract speaker notes (skip title slide)
    for slide in prs.slides[1:]:  # Skip first slide (title)
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text
        scripts.append(notes_text)
    
    # Generate voiceovers
    generator = VoiceGenerator(use_openai=True)
    audio_files = generator.generate_voiceover_for_scripts(scripts, output_dir)
    
    return audio_files


if __name__ == "__main__":
    """
    Test voiceover generation with sample scripts.
    """
    
    test_scripts = [
        "Hello! Today we're going to learn about stochastic gradient descent. "
        "This is a fundamental algorithm in machine learning.",
        
        "Let's start with the basics. Gradient descent is an optimization algorithm "
        "that finds the minimum of a function by following the negative gradient."
    ]
    
    print("🧪 Testing voiceover generation...\n")
    
    generator = VoiceGenerator(use_openai=True)
    audio_files = generator.generate_voiceover_for_scripts(
        test_scripts,
        output_dir="backend/output/test_voiceovers"
    )
    
    print(f"\n✅ Test complete! Generated {len(audio_files)} audio files")
    for path, duration in audio_files:
        print(f"   📄 {os.path.basename(path)} ({duration:.1f}s)")