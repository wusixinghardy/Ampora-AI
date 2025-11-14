import os
from typing import List, Union, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re

# Import teammate functions
from src.services.lecture import (
    generate_learning_objectives,
    generate_lecture_script
)

# TODO: Import Sam's function when it's ready
# from backend.src.services.visualization import generate_visualizations


class ScriptSplitter:
    """
    Utility class to intelligently split Hardy's long lecture script
    into per-slide chunks that match the number of visualizations.
    """
    
    @staticmethod
    def split_by_paragraphs(script: str, num_slides: int) -> List[str]:
        """
        Split script by paragraph breaks, then distribute evenly across slides.
        
        Args:
            script: Full lecture script as one string
            num_slides: Number of slides/images we need to match
            
        Returns:
            List of script chunks (one per slide)
        """
        # Split by double newlines (paragraph breaks)
        paragraphs = [p.strip() for p in script.split('\n\n') if p.strip()]
        
        if len(paragraphs) == 0:
            # Fallback: split by single newlines
            paragraphs = [p.strip() for p in script.split('\n') if p.strip()]
        
        if len(paragraphs) <= num_slides:
            # We have fewer paragraphs than slides - pad with empty strings
            while len(paragraphs) < num_slides:
                paragraphs.append("")
            return paragraphs[:num_slides]
        else:
            # We have more paragraphs than slides - group them
            chunks = []
            paragraphs_per_slide = len(paragraphs) // num_slides
            remainder = len(paragraphs) % num_slides
            
            idx = 0
            for i in range(num_slides):
                # Take extra paragraph for first 'remainder' slides
                take = paragraphs_per_slide + (1 if i < remainder else 0)
                chunk = "\n\n".join(paragraphs[idx:idx + take])
                chunks.append(chunk)
                idx += take
            
            return chunks
    
    @staticmethod
    def split_by_sentences(script: str, num_slides: int, sentences_per_slide: int = 4) -> List[str]:
        """
        Split script by sentences, grouping into chunks.
        
        Args:
            script: Full lecture script
            num_slides: Target number of slides
            sentences_per_slide: Approximate sentences per slide
            
        Returns:
            List of script chunks
        """
        # Simple sentence splitting (not perfect but works for most cases)
        sentences = re.split(r'(?<=[.!?])\s+', script)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        chunks = []
        sentences_to_take = max(1, len(sentences) // num_slides)
        
        for i in range(0, len(sentences), sentences_to_take):
            chunk = " ".join(sentences[i:i + sentences_to_take])
            chunks.append(chunk)
        
        # Ensure we have exactly num_slides chunks
        if len(chunks) < num_slides:
            chunks.extend([""] * (num_slides - len(chunks)))
        elif len(chunks) > num_slides:
            chunks = chunks[:num_slides]
        
        return chunks
    
    @staticmethod
    def split_intelligently(script: str, num_slides: int) -> List[str]:
        """
        Smart splitting: try paragraphs first, fall back to sentences.
        """
        chunks = ScriptSplitter.split_by_paragraphs(script, num_slides)
        
        # If chunks are very uneven, try sentence-based splitting
        chunk_lengths = [len(c) for c in chunks]
        if max(chunk_lengths) > 3 * min(chunk_lengths):
            chunks = ScriptSplitter.split_by_sentences(script, num_slides)
        
        return chunks


class PPTXGenerator:
    """
    PowerPoint generation class that integrates:
    - Hardy's lecture scripts (speaker notes)
    - Sam's visualizations (slide images)
    """
    
    def __init__(self):
        """Initialize presentation with standard settings"""
        self.prs = Presentation()
        # Set 16:9 widescreen format (standard for video)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def create_title_slide(self, title: str, subtitle: str = "AI-Generated Educational Content"):
        """
        Create opening title slide.
        
        Args:
            title: Main title (e.g., "Linear Regression")
            subtitle: Subtitle text
            
        Returns:
            The created slide object
        """
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
        
        return slide
    
    def create_content_slide(self, image_path: str, speaker_notes: str):
        """
        Create a content slide with:
        - Full-screen image/visualization
        - Script in speaker notes (for future voice-over)
        
        Args:
            image_path: Path to Sam's generated visualization
            speaker_notes: Portion of Hardy's script for this slide
            
        Returns:
            The created slide object
        """
        # Use blank layout for maximum image space
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Add image - try to fill most of the slide
        try:
            if not os.path.exists(image_path):
                print(f"⚠️  Warning: Image not found: {image_path}")
                # Add placeholder text instead
                left = Inches(1)
                top = Inches(3)
                width = Inches(8)
                height = Inches(1.5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.text = f"[Image Missing: {os.path.basename(image_path)}]"
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                # Add image - centered and sized to fit
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)  # Leave small margins
                
                pic = slide.shapes.add_picture(image_path, left, top, width=width)
                
                # Center image vertically
                pic.left = int((self.prs.slide_width - pic.width) / 2)
                pic.top = int((self.prs.slide_height - pic.height) / 2)
                
        except Exception as e:
            print(f"❌ Error adding image {image_path}: {e}")
        
        # Add script to speaker notes (for future voice-over generation)
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = speaker_notes
        
        return slide
    
    def generate_presentation(
        self,
        topic: str,
        lecture_script: str,
        visualization_paths: List[str],
        output_path: str = "output/generated_presentation.pptx"
    ) -> str:
        """
        Main function to generate complete presentation.
        
        Args:
            topic: Lecture topic (e.g., "Linear Regression")
            lecture_script: Full script from Hardy's generate_lecture_script()
            visualization_paths: List of image paths from Sam's function
            output_path: Where to save the .pptx file
            
        Returns:
            Path to generated presentation
        """
        
        num_slides = len(visualization_paths)
        print(f"📊 Generating presentation: {topic}")
        print(f"   Slides to create: {num_slides}")
        print(f"   Script length: {len(lecture_script)} characters")
        
        # Step 1: Split lecture script into per-slide chunks
        print("   Splitting script into per-slide chunks...")
        script_chunks = ScriptSplitter.split_intelligently(lecture_script, num_slides)
        
        if len(script_chunks) != num_slides:
            print(f"⚠️  Warning: Script chunks ({len(script_chunks)}) don't match slides ({num_slides})")
        
        # Step 2: Create title slide
        self.create_title_slide(
            title=f"Ampora: {topic}",
            subtitle="AI-Generated Educational Content"
        )
        
        # Step 3: Create content slides (image + speaker notes)
        for i, (img_path, script_chunk) in enumerate(zip(visualization_paths, script_chunks), 1):
            print(f"   Creating slide {i}/{num_slides}...")
            self.create_content_slide(img_path, script_chunk)
        
        # Step 4: Save presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        
        print(f"✅ Presentation saved: {output_path}")
        return output_path


def mock_sam_visualization(num_slides: int = 5) -> List[str]:
    """
    TEMPORARY MOCK FUNCTION for Sam's visualization.py
    
    This creates placeholder images so Robert can test pptx.py
    before Sam finishes his part.
    
    TODO: Replace this with actual call to Sam's generate_visualizations()
    
    Args:
        num_slides: Number of placeholder images to create
        
    Returns:
        List of paths to placeholder images
    """
    from PIL import Image, ImageDraw, ImageFont
    
    output_dir = "backend/output/mock_visualizations"
    os.makedirs(output_dir, exist_ok=True)
    
    image_paths = []
    
    for i in range(num_slides):
        # Create a simple placeholder image
        img = Image.new('RGB', (1920, 1080), color=(240, 240, 250))
        draw = ImageDraw.Draw(img)
        
        # Add text
        text = f"Slide {i+1}\n(Sam's visualization will go here)"
        
        # Try to use a font, fall back to default if not available
        try:
            font = ImageFont.truetype("arial.ttf", 60)
        except:
            font = ImageFont.load_default()
        
        # Calculate text position (centered)
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        position = ((1920 - text_width) / 2, (1080 - text_height) / 2)
        
        draw.text(position, text, fill=(100, 100, 100), font=font)
        
        # Save
        img_path = os.path.join(output_dir, f"slide_{i+1}.png")
        img.save(img_path)
        image_paths.append(img_path)
    
    return image_paths


def create_presentation_from_topic(
    topic: str,
    output_path: Optional[str] = None,
    # In create_presentation_from_topic(), change:
    # use_mock_visualizations=False

    # And make sure Sam's function signature matches:
    # def generate_visualizations(lecture_script: str) -> List[str]:
    #     """Returns list of image file paths"""
    use_mock_visualizations: bool = True
) -> str:
    """
    High-level function that orchestrates the full pipeline:
    1. Call Hardy's functions to generate learning objectives and script
    2. Call Sam's function to generate visualizations (or use mock for testing)
    3. Generate PowerPoint with integrated content
    
    Args:
        topic: Topic to teach (e.g., "Linear Regression", "Convolutional Neural Networks")
        output_path: Where to save the presentation (auto-generated if None)
        use_mock_visualizations: If True, use placeholder images until Sam's function is ready
        
    Returns:
        Path to generated .pptx file
    """
    
    print("="*70)
    print(f"AMPORA PRESENTATION GENERATOR - {topic}")
    print("="*70)
    
    # Step 1: Generate learning objectives (Hardy's function)
    print("\n📚 Step 1: Generating learning objectives...")
    try:
        objectives = generate_learning_objectives(topic)
        print(f"   ✅ Generated {len(objectives)} learning objectives")
        for i, obj in enumerate(objectives, 1):
            print(f"      {i}. {obj[:80]}...")
    except Exception as e:
        print(f"   ❌ Error generating objectives: {e}")
        raise
    
    # Step 2: Generate lecture script (Hardy's function)
    print("\n📝 Step 2: Generating lecture script...")
    try:
        lecture_script = generate_lecture_script(objectives)
        print(f"   ✅ Generated script ({len(lecture_script)} characters)")
        print(f"   Preview: {lecture_script[:200]}...")
    except Exception as e:
        print(f"   ❌ Error generating script: {e}")
        raise
    
    # Step 3: Generate visualizations (Sam's function - or mock for now)
    print("\n🎨 Step 3: Generating visualizations...")
    try:
        if use_mock_visualizations:
            print("   ⚠️  Using mock visualizations (Sam's function not ready yet)")
            # Estimate number of slides based on script length
            # ~200 words per slide, ~5 chars per word = ~1000 chars per slide
            estimated_slides = max(3, min(10, len(lecture_script) // 1000))
            visualization_paths = mock_sam_visualization(estimated_slides)
        else:
            # TODO: Uncomment when Sam's function is ready
            # from backend.src.services.visualization import generate_visualizations
            # visualization_paths = generate_visualizations(lecture_script)
            raise NotImplementedError("Sam's generate_visualizations() not implemented yet")
        
        print(f"   ✅ Generated {len(visualization_paths)} visualizations")
        for path in visualization_paths:
            print(f"      - {path}")
    except Exception as e:
        print(f"   ❌ Error generating visualizations: {e}")
        raise
    
    # Step 4: Create PowerPoint (Robert's function - this file!)
    print("\n📊 Step 4: Creating PowerPoint presentation...")
    try:
        generator = PPTXGenerator()
        
        if output_path is None:
            output_path = f"backend/output/{topic.replace(' ', '_')}_presentation.pptx"
        
        final_path = generator.generate_presentation(
            topic=topic,
            lecture_script=lecture_script,
            visualization_paths=visualization_paths,
            output_path=output_path
        )
    except Exception as e:
        print(f"   ❌ Error creating PowerPoint: {e}")
        raise
    
    # Success!
    print("\n" + "="*70)
    print("✅ PIPELINE COMPLETE!")
    print("="*70)
    print(f"📄 Presentation: {final_path}")
    print(f"📊 Slides: {len(visualization_paths) + 1} (including title)")
    print(f"📝 Script: {len(lecture_script)} characters")
    print("\nNext steps:")
    print("1. Review the generated presentation")
    print("2. Test with different topics")
    print("3. Wait for Sam to implement visualization.py, then set use_mock_visualizations=False")
    print("4. Eventually: Add voice-over generation and MP4 export")
    print("="*70 + "\n")
    
    return final_path


# Test/Example usage
if __name__ == "__main__":
    """
    Test the pipeline with a sample topic.
    Run this to verify everything works before integrating with teammates' code.
    """
    
    # Test with a simple topic
    test_topic = "Linear Regression"
    
    print("🧪 Running test with topic:", test_topic)
    print("⚠️  Note: Using mock visualizations until Sam's function is ready\n")
    
    try:
        output_file = create_presentation_from_topic(
            topic=test_topic,
            use_mock_visualizations=True  # Set to False when Sam's function is ready
        )
        
        print(f"\n✅ Test successful! Open the file to review: {output_file}")
        
    except Exception as e:
        print(f"\n❌ Test failed: {e}")
        import traceback
        traceback.print_exc()